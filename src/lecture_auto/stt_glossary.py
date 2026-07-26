from __future__ import annotations

import re
from collections import Counter
from pathlib import Path
from typing import Iterable


MAX_GLOSSARY_TERMS = 64
MAX_HOTWORDS_CHARACTERS = 1000

_ENGLISH_TERM = re.compile(
    r"\b(?:[A-Z][A-Za-z0-9+.#-]{1,}|[A-Z]{2,}(?:-[A-Z0-9]+)*|"
    r"[A-Za-z][A-Za-z0-9+.#-]+(?:\s+[A-Za-z][A-Za-z0-9+.#-]+){1,3})\b"
)
_KOREAN_TERM = re.compile(r"(?<![가-힣])(?:[가-힣]{2,12})(?![가-힣])")
_STOPWORDS = {
    "그리고",
    "그러나",
    "그래서",
    "이것은",
    "입니다",
    "합니다",
    "대한",
    "있는",
    "없는",
    "the",
    "and",
    "this",
    "that",
    "with",
    "from",
    "lecture",
}


def extract_glossary_terms(
    *,
    title: str | None = None,
    course: str | None = None,
    material_path: str | Path | None = None,
    max_terms: int = MAX_GLOSSARY_TERMS,
    max_characters: int = MAX_HOTWORDS_CHARACTERS,
) -> list[str]:
    seed_sources = [
        value.strip() for value in (title, course) if value and value.strip()
    ]
    terms = bounded_glossary(
        seed_sources,
        max_terms=min(16, max_terms),
        max_characters=max_characters,
    )
    if material_path:
        material_terms = bounded_glossary(
            [_extract_material_text(Path(material_path))[:200_000]],
            max_terms=max_terms,
            max_characters=max_characters,
            minimum_frequency=2,
        )
        for term in material_terms:
            if term.casefold() not in {value.casefold() for value in terms}:
                terms.append(term)
    return _limit_terms(
        terms,
        max_terms=max_terms,
        max_characters=max_characters,
    )


def bounded_glossary(
    sources: Iterable[str],
    *,
    max_terms: int = MAX_GLOSSARY_TERMS,
    max_characters: int = MAX_HOTWORDS_CHARACTERS,
    minimum_frequency: int = 1,
) -> list[str]:
    if max_terms < 1 or max_characters < 1:
        return []

    candidates: list[str] = []
    for source in sources:
        text = " ".join(str(source).split())
        candidates.extend(match.group(0).strip() for match in _ENGLISH_TERM.finditer(text))
        candidates.extend(match.group(0).strip() for match in _KOREAN_TERM.finditer(text))

    counts = Counter(term.casefold() for term in candidates if _valid_term(term))
    first_form: dict[str, str] = {}
    for term in candidates:
        folded = term.casefold()
        if _valid_term(term):
            first_form.setdefault(folded, term)

    ordered = sorted(
        (term for term in first_form if counts[term] >= minimum_frequency),
        key=lambda term: (-counts[term], -len(first_form[term]), term),
    )
    return _limit_terms(
        (first_form[folded] for folded in ordered),
        max_terms=max_terms,
        max_characters=max_characters,
    )


def _limit_terms(
    terms: Iterable[str],
    *,
    max_terms: int,
    max_characters: int,
) -> list[str]:
    selected: list[str] = []
    used = 0
    for term in terms:
        added = len(term) + (1 if selected else 0)
        if used + added > max_characters:
            continue
        selected.append(term)
        used += added
        if len(selected) >= max_terms:
            break
    return selected


def merge_hotwords(configured: str | None, glossary: Iterable[str]) -> str | None:
    sources = [configured or "", *glossary]
    terms = bounded_glossary(
        sources,
        max_terms=MAX_GLOSSARY_TERMS,
        max_characters=MAX_HOTWORDS_CHARACTERS,
    )
    return " ".join(terms) or None


def _valid_term(term: str) -> bool:
    normalized = term.strip(" .,:;()[]{}")
    if len(normalized) < 2 or normalized.casefold() in _STOPWORDS:
        return False
    if normalized.isdigit():
        return False
    return True


def _extract_material_text(path: Path) -> str:
    if not path.is_file():
        return ""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader

            return "\n".join((page.extract_text() or "") for page in PdfReader(str(path)).pages)
        if suffix == ".pptx":
            from pptx import Presentation

            return "\n".join(
                shape.text
                for slide in Presentation(str(path)).slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
        if suffix in {".md", ".txt"}:
            return path.read_text(encoding="utf-8", errors="ignore")
    except (OSError, ValueError):
        return ""
    return ""
